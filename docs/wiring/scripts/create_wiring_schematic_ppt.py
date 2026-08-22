#!/usr/bin/env python3
"""
Create PowerPoint wiring schematic presentation
- First slide: Summary of all nodes with board counts
- One slide per control node (C1-C13) with separate tables for each board
"""
import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from collections import defaultdict
import os
import re

script_dir = os.path.dirname(os.path.abspath(__file__))
wiring_dir = os.path.normpath(os.path.join(script_dir, '..'))
# Find latest inventory file (skip imported/ snapshots)
import glob
inventory_files = sorted(
    glob.glob(os.path.join(wiring_dir, 'LCOS_Layout_Inventory_v*.xlsx')),
    reverse=True,
)
inventory_file = inventory_files[0] if inventory_files else os.path.join(wiring_dir, 'LCOS_Layout_Inventory_v85.xlsx')
output_ppt = os.path.join(wiring_dir, 'Wiring_Schematic.pptx')

print("=" * 80)
print("CREATING WIRING SCHEMATIC POWERPOINT PRESENTATION")
print("=" * 80)

# Load Excel file
print(f"Using inventory file: {os.path.basename(inventory_file)}")
wb = openpyxl.load_workbook(inventory_file)
ws_nodes = wb['Nodes']
ws_dnou8 = wb['DNOU8']
ws_dnin8 = wb['DNIN8']
ws_bs = wb['BlockSensors']
ws_ts = wb['TurnoutSummary'] if 'TurnoutSummary' in wb.sheetnames else None

# Get headers
nodes_headers = [cell.value for cell in ws_nodes[1]]
dnou8_headers = [cell.value for cell in ws_dnou8[1]]
dnin8_headers = [cell.value for cell in ws_dnin8[1]]
bs_headers = [cell.value for cell in ws_bs[1]]

# Column indices
node_id_col = nodes_headers.index('Node ID')
location_col = nodes_headers.index('Location') if 'Location' in nodes_headers else None
address_col = nodes_headers.index('Address') if 'Address' in nodes_headers else None
boards_12v_col = nodes_headers.index('12V Boards') if '12V Boards' in nodes_headers else None
boards_5v_col = nodes_headers.index('5V Boards') if '5V Boards' in nodes_headers else None
boards_input_col = nodes_headers.index('Input Boards') if 'Input Boards' in nodes_headers else None
num_blocks_col = nodes_headers.index('Num Blocks') if 'Num Blocks' in nodes_headers else None

dnou8_node_col = dnou8_headers.index('Parent Node ID')
dnou8_port_col = dnou8_headers.index('Output Port ID')
dnou8_device_col = dnou8_headers.index('Connected Device')

dnin8_node_col = dnin8_headers.index('Parent Node ID')
dnin8_port_col = dnin8_headers.index('Input Port ID')
dnin8_device_col = dnin8_headers.index('Connected Device')

bs_node_col = bs_headers.index('Node ID')
bs_port_col = bs_headers.index('Port ID')
bs_block_name_col = bs_headers.index('Block Section Name')

# Collect node summary data
print("\nCollecting node data...")
node_summaries = {}
node_data = {}

for row_idx, row in enumerate(ws_nodes.iter_rows(min_row=2), start=2):
    node_id = ws_nodes.cell(row_idx, node_id_col + 1).value
    if not node_id or not str(node_id).startswith('C'):
        continue
    
    location = ws_nodes.cell(row_idx, location_col + 1).value if location_col else None
    address = ws_nodes.cell(row_idx, address_col + 1).value if address_col else None
    boards_12v = ws_nodes.cell(row_idx, boards_12v_col + 1).value if boards_12v_col else 0
    boards_5v = ws_nodes.cell(row_idx, boards_5v_col + 1).value if boards_5v_col else 0
    boards_input = ws_nodes.cell(row_idx, boards_input_col + 1).value if boards_input_col else 0
    num_blocks = ws_nodes.cell(row_idx, num_blocks_col + 1).value if num_blocks_col else 0
    
    # Count input boards from actual DNIN8 data (will be updated after collecting DNIN8)
    node_summaries[node_id] = {
        'location': location or 'N/A',
        'address': address or 'N/A',
        'boards_12v': int(boards_12v) if boards_12v else 0,
        'boards_5v': int(boards_5v) if boards_5v else 0,
        'boards_input': 0,  # Will be calculated from DNIN8 data
        'num_blocks': int(num_blocks) if num_blocks else 0
    }
    
    node_data[node_id] = {
        'ou_boards': defaultdict(dict),  # board_id -> {port_num: device}
        'in_boards': defaultdict(dict),  # board_id -> {port_num: device}
        'blocks': []
    }

# Helper function to extract board ID from port
def get_board_id(port_str):
    """Extract board ID from port string (e.g., 'C3-OU1-1' -> 'C3-OU1')"""
    if not port_str:
        return None
    parts = str(port_str).split('-')
    if len(parts) >= 3:
        return f"{parts[0]}-{parts[1]}"
    return None

# Helper function to extract port number from port string
def get_port_num(port_str):
    """Extract port number from port string (e.g., 'C3-OU1-1' -> 1)"""
    if not port_str:
        return None
    parts = str(port_str).split('-')
    if len(parts) >= 3:
        try:
            return int(parts[2])
        except:
            return None
    return None

# Collect DNOU8 data grouped by board
print("Collecting DNOU8 data...")
for row in ws_dnou8.iter_rows(min_row=2, values_only=True):
    node = row[dnou8_node_col]
    port = row[dnou8_port_col]
    device = row[dnou8_device_col]
    
    if node and port and node in node_data:
        board_id = get_board_id(port)
        port_num = get_port_num(port)
        
        if board_id and port_num:
            device_str = str(device).strip() if device else ''
            node_data[node]['ou_boards'][board_id][port_num] = device_str

# Collect DNIN8 data grouped by board
print("Collecting DNIN8 data...")
for row in ws_dnin8.iter_rows(min_row=2, values_only=True):
    node = row[dnin8_node_col]
    port = row[dnin8_port_col]
    device = row[dnin8_device_col]
    
    if node and port and node in node_data:
        board_id = get_board_id(port)
        port_num = get_port_num(port)
        
        if board_id and port_num:
            device_str = str(device).strip() if device else ''
            node_data[node]['in_boards'][board_id][port_num] = device_str

# Update input board counts from actual DNIN8 data
print("Calculating input board counts...")
for node_id in node_summaries:
    if node_id in node_data:
        input_boards_count = len(node_data[node_id]['in_boards'])
        node_summaries[node_id]['boards_input'] = input_boards_count

# Collect BlockSensors data
print("Collecting BlockSensors data...")
for row in ws_bs.iter_rows(min_row=2, values_only=True):
    node = row[bs_node_col]
    port = row[bs_port_col]
    block_name = row[bs_block_name_col]
    
    if node and port and block_name and str(block_name).strip():
        if node in node_data:
            node_data[node]['blocks'].append({
                'port': str(port),
                'block': str(block_name).strip()
            })

print(f"Collected data for {len(node_data)} nodes")

# Create PowerPoint presentation
print("\nCreating PowerPoint presentation...")
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Slide 1: Summary of all nodes
print("Creating summary slide...")
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
title_frame = title_box.text_frame
title_frame.text = "Wiring Schematic - Node Summary"
title_para = title_frame.paragraphs[0]
title_para.font.size = Pt(32)
title_para.font.bold = True

# Create summary table
sorted_nodes = sorted([n for n in node_summaries.keys() if str(n).startswith('C')], 
                     key=lambda x: int(x[1:]) if x[1:].isdigit() else 999)
rows = len(sorted_nodes) + 1
cols = 7
left = Inches(0.5)
top = Inches(1.5)
width = Inches(9)
height = Inches(5.5)

table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Header row
table.cell(0, 0).text = "Node ID"
table.cell(0, 1).text = "Location"
table.cell(0, 2).text = "Radio Address"
table.cell(0, 3).text = "12V Boards"
table.cell(0, 4).text = "5V Boards"
table.cell(0, 5).text = "Input Boards"
table.cell(0, 6).text = "Num Blocks"

# Format header
for col in range(cols):
    cell = table.cell(0, col)
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
    para = cell.text_frame.paragraphs[0]
    para.font.color.rgb = RGBColor(255, 255, 255)
    para.font.bold = True
    para.font.size = Pt(12)

# Add node data
for idx, node_id in enumerate(sorted_nodes, 1):
    summary = node_summaries[node_id]
    
    table.cell(idx, 0).text = str(node_id)
    table.cell(idx, 1).text = str(summary['location'])
    table.cell(idx, 2).text = str(summary['address'])
    table.cell(idx, 3).text = str(summary['boards_12v'])
    table.cell(idx, 4).text = str(summary['boards_5v'])
    table.cell(idx, 5).text = str(summary['boards_input'])
    table.cell(idx, 6).text = str(summary['num_blocks'])
    
    # Format data rows
    for col in range(cols):
        cell = table.cell(idx, col)
        para = cell.text_frame.paragraphs[0]
        para.font.size = Pt(10)
        if idx % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

# Create Turnout Summary slide
if ws_ts:
    print("Creating turnout summary slide...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = "Turnout Summary"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    
    # Get TurnoutSummary headers
    ts_headers = [cell.value for cell in ws_ts[1]]
    ts_turnout_col = ts_headers.index('Turnout') if 'Turnout' in ts_headers else None
    ts_node_col = ts_headers.index('Parent Node') if 'Parent Node' in ts_headers else None
    ts_location_col = ts_headers.index('Location/Area') if 'Location/Area' in ts_headers else None
    ts_button_col = ts_headers.index('Button') if 'Button' in ts_headers else None
    ts_fb_n_col = ts_headers.index('Feedback N') if 'Feedback N' in ts_headers else None
    ts_fb_r_col = ts_headers.index('Feedback R') if 'Feedback R' in ts_headers else None
    ts_entry_signal_col = ts_headers.index('Entry Signal') if 'Entry Signal' in ts_headers else None
    ts_normal_signal_col = ts_headers.index('Normal Exit Signal') if 'Normal Exit Signal' in ts_headers else None
    ts_reverse_signal_col = ts_headers.index('Reverse Exit Signal') if 'Reverse Exit Signal' in ts_headers else None
    
    # Collect turnout data
    turnout_data = []
    for row_idx, row in enumerate(ws_ts.iter_rows(min_row=2), start=2):
        turnout = ws_ts.cell(row_idx, ts_turnout_col + 1).value if ts_turnout_col is not None else None
        if not turnout:
            continue
        
        node = ws_ts.cell(row_idx, ts_node_col + 1).value if ts_node_col is not None else None
        location = ws_ts.cell(row_idx, ts_location_col + 1).value if ts_location_col is not None else None
        button = ws_ts.cell(row_idx, ts_button_col + 1).value if ts_button_col is not None else None
        fb_n = ws_ts.cell(row_idx, ts_fb_n_col + 1).value if ts_fb_n_col is not None else None
        fb_r = ws_ts.cell(row_idx, ts_fb_r_col + 1).value if ts_fb_r_col is not None else None
        
        # Collect associated signals
        entry_signal = ws_ts.cell(row_idx, ts_entry_signal_col + 1).value if ts_entry_signal_col is not None else None
        normal_signal = ws_ts.cell(row_idx, ts_normal_signal_col + 1).value if ts_normal_signal_col is not None else None
        reverse_signal = ws_ts.cell(row_idx, ts_reverse_signal_col + 1).value if ts_reverse_signal_col is not None else None
        
        # Combine feedback ports
        feedback_ports = []
        if fb_n:
            feedback_ports.append(f"N:{fb_n}")
        if fb_r:
            feedback_ports.append(f"R:{fb_r}")
        feedback_str = ", ".join(feedback_ports) if feedback_ports else ''
        
        # Combine signals (just names, comma-separated)
        signals = []
        if entry_signal:
            signals.append(str(entry_signal))
        if normal_signal:
            signals.append(str(normal_signal))
        if reverse_signal:
            signals.append(str(reverse_signal))
        signals_str = ", ".join(signals) if signals else ''
        
        turnout_data.append({
            'turnout': str(turnout),
            'node': str(node) if node else '',
            'button': str(button) if button else '',
            'feedback': feedback_str,
            'signals': signals_str
        })
    
    # Sort by control node (C1-C13), then turnout
    def get_node_number(node_str):
        if node_str and node_str.startswith('C'):
            try:
                return int(node_str[1:])
            except:
                return 999
        return 999
    
    turnout_data.sort(key=lambda x: (get_node_number(x['node']), x['turnout']))
    
    # Split into two columns
    mid_point = (len(turnout_data) + 1) // 2
    left_column_data = turnout_data[:mid_point]
    right_column_data = turnout_data[mid_point:]
    
    # Calculate table dimensions
    left_rows = len(left_column_data) + 1
    right_rows = len(right_column_data) + 1
    cols = 5  # Removed location column
    # Table width (sum of column widths: 0.7+0.5+0.9+1.4+1.0 = 4.5")
    table_width = Inches(4.5)
    table_height = Inches(5.5)
    
    # Use specific coordinates provided by user
    left_table_x = Inches(0.32)
    right_table_x = Inches(5.18)
    top = Inches(1.41)
    
    # Create left table
    left_table = slide.shapes.add_table(left_rows, cols, left_table_x, top, table_width, table_height).table
    
    # Header row for left table
    left_table.cell(0, 0).text = "Turnout"
    left_table.cell(0, 1).text = "Node"
    left_table.cell(0, 2).text = "Button"
    left_table.cell(0, 3).text = "Feedback"
    left_table.cell(0, 4).text = "Signals"
    
    # Format left header
    for col in range(cols):
        cell = left_table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.bold = True
        para.font.size = Pt(9)
    
    # Set column widths (in EMU: 1 inch = 914400 EMU)
    # Turnout: narrow, Node: narrow, Button: narrow, Feedback: wider, Signals: narrower
    left_table.columns[0].width = Inches(0.7)  # Turnout
    left_table.columns[1].width = Inches(0.5)  # Node
    left_table.columns[2].width = Inches(0.9)  # Button
    left_table.columns[3].width = Inches(1.4)  # Feedback (wider)
    left_table.columns[4].width = Inches(1.0)  # Signals (narrower)
    
    # Add left data
    for idx, data in enumerate(left_column_data, 1):
        left_table.cell(idx, 0).text = data['turnout']
        left_table.cell(idx, 1).text = data['node']
        left_table.cell(idx, 2).text = data['button']
        left_table.cell(idx, 3).text = data['feedback']
        left_table.cell(idx, 4).text = data['signals']
        
        # Format data rows
        for col in range(cols):
            cell = left_table.cell(idx, col)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(8)
            if idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    # Create right table
    right_table = slide.shapes.add_table(right_rows, cols, right_table_x, top, table_width, table_height).table
    
    # Header row for right table
    right_table.cell(0, 0).text = "Turnout"
    right_table.cell(0, 1).text = "Node"
    right_table.cell(0, 2).text = "Button"
    right_table.cell(0, 3).text = "Feedback"
    right_table.cell(0, 4).text = "Signals"
    
    # Format right header
    for col in range(cols):
        cell = right_table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        para = cell.text_frame.paragraphs[0]
        para.font.color.rgb = RGBColor(255, 255, 255)
        para.font.bold = True
        para.font.size = Pt(9)
    
    # Set column widths (same as left table)
    right_table.columns[0].width = Inches(0.7)  # Turnout
    right_table.columns[1].width = Inches(0.5)  # Node
    right_table.columns[2].width = Inches(0.9)  # Button
    right_table.columns[3].width = Inches(1.4)  # Feedback (wider)
    right_table.columns[4].width = Inches(1.0)  # Signals (narrower)
    
    # Add right data
    for idx, data in enumerate(right_column_data, 1):
        right_table.cell(idx, 0).text = data['turnout']
        right_table.cell(idx, 1).text = data['node']
        right_table.cell(idx, 2).text = data['button']
        right_table.cell(idx, 3).text = data['feedback']
        right_table.cell(idx, 4).text = data['signals']
        
        # Format data rows
        for col in range(cols):
            cell = right_table.cell(idx, col)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(8)
            if idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)

# Create individual slides for each node
print("Creating individual node slides...")
prev_location = None
for node_id in sorted_nodes:
    print(f"  Creating slide for {node_id}...")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Title with location
    summary = node_summaries[node_id]
    current_location = summary['location']
    
    # Add "(cont'd)" if same location as previous slide
    if prev_location and current_location == prev_location:
        title_text = f"Node {node_id} - {current_location} (cont'd)"
    else:
        title_text = f"Node {node_id} - {current_location}"
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = title_text
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    
    prev_location = current_location
    
    data = node_data[node_id]
    
    # Define 6 sections on the slide (2 columns x 3 rows)
    # Section layout:
    # [0] [1]
    # [2] [3]
    # [4] [5]
    
    section_width = Inches(4.5)
    section_height = Inches(1.7)  # Adjusted height
    section_spacing_x = Inches(0.2)
    section_spacing_y = Inches(0.15)
    start_x = Inches(0.5)
    start_y = Inches(1.0)  # Start higher since we removed the location line
    
    # User-specified positions
    top_y = start_y  # Top row
    middle_y = Inches(3.09)  # Middle row
    bottom_y = Inches(5.27)  # Bottom row
    
    section_positions = [
        (start_x, top_y),  # Section 0: top-left
        (start_x + section_width + section_spacing_x, top_y),  # Section 1: top-right
        (start_x, middle_y),  # Section 2: middle-left
        (start_x + section_width + section_spacing_x, middle_y),  # Section 3: middle-right
        (start_x, bottom_y),  # Section 4: bottom-left
        (start_x + section_width + section_spacing_x, bottom_y),  # Section 5: bottom-right
    ]
    
    table_width = section_width
    table_height = section_height
    
    # Fill left side first (sections 0, 2, 4), then right side (sections 1, 3, 5)
    left_sections = [0, 2, 4]  # top-left, middle-left, bottom-left
    right_sections = [1, 3, 5]  # top-right, middle-right, bottom-right
    
    left_index = 0
    right_index = 0
    
    # Create tables for each OU board (left side)
    ou_boards = sorted(data['ou_boards'].keys())
    for board_id in ou_boards:
        if left_index >= len(left_sections):
            break  # No more left sections
        
        section_index = left_sections[left_index]
        current_x, current_y = section_positions[section_index]
        left_index += 1
        
        # Create table with 9 rows (board title + 8 ports, no column headers)
        ou_table = slide.shapes.add_table(9, 2, current_x, current_y, table_width, section_height).table
        
        # Board title row (merged)
        title_cell = ou_table.cell(0, 0)
        title_cell.text = board_id
        title_cell.merge(ou_table.cell(0, 1))
        title_para = title_cell.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(10)
        title_para.alignment = PP_ALIGN.CENTER
        title_cell.fill.solid()
        title_cell.fill.fore_color.rgb = RGBColor(68, 114, 196)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add 8 port rows (no header row)
        board_data = data['ou_boards'][board_id]
        for port_num in range(1, 9):
            row_idx = port_num  # Row 1-8 (row 0 is board title)
            port_str = f"{board_id}-{port_num}"
            device = board_data.get(port_num, '')
            
            ou_table.cell(row_idx, 0).text = port_str
            ou_table.cell(row_idx, 1).text = device if device else ''
            
            for col in range(2):
                cell = ou_table.cell(row_idx, col)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(8)  # Normal font size
                cell.vertical_anchor = 1  # Middle alignment
                # Reduce paragraph spacing
                para.space_after = Pt(0)
                para.space_before = Pt(0)
                para.line_spacing = 0.9  # Tighter line spacing
    
    # Create tables for each IN board (right side)
    in_boards = sorted(data['in_boards'].keys())
    for board_id in in_boards:
        if right_index >= len(right_sections):
            break  # No more right sections
        
        section_index = right_sections[right_index]
        current_x, current_y = section_positions[section_index]
        right_index += 1
        
        # Create table with 9 rows (board title + 8 ports, no column headers)
        in_table = slide.shapes.add_table(9, 2, current_x, current_y, table_width, section_height).table
        
        # Board title row (merged)
        title_cell = in_table.cell(0, 0)
        title_cell.text = board_id
        title_cell.merge(in_table.cell(0, 1))
        title_para = title_cell.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(10)
        title_para.alignment = PP_ALIGN.CENTER
        title_cell.fill.solid()
        title_cell.fill.fore_color.rgb = RGBColor(112, 173, 71)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add 8 port rows (no header row)
        board_data = data['in_boards'][board_id]
        for port_num in range(1, 9):
            row_idx = port_num  # Row 1-8 (row 0 is board title)
            port_str = f"{board_id}-{port_num}"
            device = board_data.get(port_num, '')
            
            in_table.cell(row_idx, 0).text = port_str
            in_table.cell(row_idx, 1).text = device if device else ''
            
            for col in range(2):
                cell = in_table.cell(row_idx, col)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(8)  # Normal font size
                cell.vertical_anchor = 1  # Middle alignment
                # Reduce paragraph spacing
                para.space_after = Pt(0)
                para.space_before = Pt(0)
                para.line_spacing = 0.9  # Tighter line spacing
    
    # Blocks table (use remaining right section if available)
    if data['blocks'] and right_index < len(right_sections):
        # Limit to 8 rows to match board tables (or use actual count if less)
        blocks_data = sorted(data['blocks'], key=lambda x: x['port'])[:8]
        blocks_rows = len(blocks_data) + 1  # +1 for title row to match board tables
        blocks_cols = 2
        section_index = right_sections[right_index]
        current_x, current_y = section_positions[section_index]
        # Use same fixed height as board tables
        blocks_table = slide.shapes.add_table(blocks_rows, blocks_cols,
                                              current_x, current_y,
                                              table_width, section_height).table
        
        # Add title row (merged) to match board table structure
        title_cell = blocks_table.cell(0, 0)
        title_cell.text = f"{node_id} Blocks"
        title_cell.merge(blocks_table.cell(0, 1))
        title_para = title_cell.text_frame.paragraphs[0]
        title_para.font.bold = True
        title_para.font.size = Pt(10)
        title_para.alignment = PP_ALIGN.CENTER
        title_cell.fill.solid()
        title_cell.fill.fore_color.rgb = RGBColor(237, 125, 49)
        title_para.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add data (starting at row 1, matching board tables)
        for idx, entry in enumerate(blocks_data):
            row_idx = idx + 1  # Row 1-8 (row 0 is title)
            blocks_table.cell(row_idx, 0).text = entry['port']
            blocks_table.cell(row_idx, 1).text = entry['block']
            for col in range(blocks_cols):
                cell = blocks_table.cell(row_idx, col)
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(8)  # Same font size as board tables
                cell.vertical_anchor = 1  # Middle alignment
                # Reduce paragraph spacing
                para.space_after = Pt(0)
                para.space_before = Pt(0)
                para.line_spacing = 0.9  # Tighter line spacing

# Save presentation
prs.save(output_ppt)
print(f"\n✓ Saved presentation to {output_ppt}")
print(f"✓ Created {len(prs.slides)} slides total")
print(f"  - 1 summary slide")
print(f"  - {len(sorted_nodes)} individual node slides")

print("\n" + "=" * 80)
print("COMPLETE")
print("=" * 80)
